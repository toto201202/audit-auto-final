from flask import Flask, request, send_file
from pptx import Presentation
import io

app = Flask(__name__)

@app.route("/")
def index():
    return "Hello from Flask API"

@app.route("/generate", methods=["POST"])
def generate_pptx():
    data = request.json
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = f"Audit Report: {data.get('site_name', 'No site name')}"
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return send_file(output, download_name="audit_report.pptx", as_attachment=True)

if __name__ == "__main__":
    app.run(debug=True)